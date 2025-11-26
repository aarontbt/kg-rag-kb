from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from typing import List, Dict, Any
import logging

logger = logging.getLogger(__name__)

class PPTXProcessor:
    """Process PowerPoint presentations to extract text and metadata"""
    
    def __init__(self):
        pass
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a PowerPoint file and return extracted text with metadata
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of dictionaries containing extracted text and metadata
        """
        file_path = Path(file_path)
        chunks = []
        
        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = self._extract_slide_text(slide)
                
                if slide_text:
                    chunks.append({
                        'text': slide_text.strip(),
                        'metadata': {
                            'source': str(file_path),
                            'file_type': 'pptx',
                            'slide': slide_num + 1,
                            'total_slides': total_slides
                        }
                    })
            
            logger.info(f"Successfully processed PowerPoint: {file_path}")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
            return []
    
    def _extract_slide_text(self, slide) -> str:
        """
        Extract all text content from a slide
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Extracted text from the slide
        """
        text_parts = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            shape_text = self._extract_shape_text(shape)
            if shape_text:
                text_parts.append(shape_text)
        
        # Extract text from tables
        for shape in slide.shapes:
            if shape.has_table:
                table_text = self._extract_table_text(shape.table)
                if table_text:
                    text_parts.append(table_text)
        
        # Join all text with newlines
        return '\n'.join(text_parts)
    
    def _extract_shape_text(self, shape) -> str:
        """
        Extract text from a shape
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            Extracted text from the shape
        """
        try:
            # Check if shape has text frame
            if hasattr(shape, "text_frame") and shape.text_frame:
                return shape.text_frame.text
            
            # Check if shape is a picture (might have alt text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if hasattr(shape, "alt_text"):
                    return shape.alt_text
            
            return ""
            
        except Exception as e:
            logger.warning(f"Error extracting text from shape: {str(e)}")
            return ""
    
    def _extract_table_text(self, table) -> str:
        """
        Extract text from a table
        
        Args:
            table: PowerPoint table object
            
        Returns:
            Extracted text from the table in CSV-like format
        """
        try:
            rows_text = []
            
            for row in table.rows:
                cells_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    cells_text.append(cell_text)
                
                if any(cells_text):  # Skip empty rows
                    rows_text.append(' | '.join(cells_text))
            
            return '\n'.join(rows_text)
            
        except Exception as e:
            logger.warning(f"Error extracting text from table: {str(e)}")
            return ""
    
    def is_supported(self, file_path: str) -> bool:
        """Check if the file is a supported PowerPoint format"""
        return Path(file_path).suffix.lower() in ['.pptx', '.ppt']
